"""
extractor.py : Convert file formats to plain text for vault ingestion.

Supported: .md/.markdown/.mdx  .txt/.text  .csv  .html/.htm  .pdf  .docx  .xlsx  .pptx  .json
Images are not supported: convert to .txt before dropping in the inbox.

All extractors return a string. Callers treat None or empty string as a failure
and route the file to _failed/.

Recursive split helpers (v0.3):
  extract_pages(path)    -> list[str]  per-page text for PDFs; [] otherwise
"""

import csv as _csv
import json
import logging
from pathlib import Path

logger = logging.getLogger("extractor")

SUPPORTED: frozenset[str] = frozenset({
    # .mdx is markdown with JSX components interleaved. The prose reads fine as
    # text and the components degrade to inert tags, which is why it goes through
    # _text unchanged: openclaw's docs skipped three real deploy guides over an
    # extension that needed no extraction logic of its own.
    ".md", ".markdown", ".mdx",
    ".txt", ".text",
    ".csv",
    ".html", ".htm",
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".json",
})


class DatalessFileError(OSError):
    """File is evicted to the cloud (iCloud, OneDrive, Dropbox) and has st_blocks==0."""
    pass


class UnreadableFileError(Exception):
    """File could not be parsed or decoded."""
    pass


def _first_byte_readable(path: Path) -> bool:
    """True when one byte comes back off the file without blocking on a download.

    The discriminator between a genuinely evicted file and a filesystem that
    simply does not report block counts: an evicted file raises (commonly
    EDEADLK, "Resource deadlock avoided", under launchd) or yields nothing,
    while an inline-extent file on btrfs/ext4 reads back normally.
    """
    try:
        with open(path, "rb") as f:
            return bool(f.read(1))
    except Exception:
        return False


def is_dataless(path: Path) -> bool:
    """Return True if a file exists on disk but its content is evicted to cloud storage."""
    try:
        st = path.stat()
        # st_blocks == 0 on a non-empty file means "evicted to the cloud" only
        # where the filesystem allocates blocks for content at all. btrfs and
        # ext4 store small files inline in metadata, and FUSE/network mounts
        # frequently report no block count — on those, a bare st_blocks test
        # grades an entire healthy tree as dataless, materialize fails to move
        # the number, and ingest reports every file skipped_dataless while
        # indexing nothing. The flag below is the authoritative macOS signal;
        # st_blocks is a fallback, so it must also fail to actually read.
        flags = getattr(st, "st_flags", 0)
        if flags & 0x40000000:  # UF_DATALESS / SF_DATALESS on BSD/macOS
            return True
        if st.st_size > 0 and getattr(st, "st_blocks", 1) == 0:
            return not _first_byte_readable(path)
    except Exception:
        pass
    return False


def materialize_if_dataless(path: Path, timeout: float = 10.0) -> bool:
    """Attempt to trigger on-demand cloud file download if dataless. Returns True if readable."""
    if not is_dataless(path):
        return True
    try:
        import shutil, subprocess
        if shutil.which("brctl"):
            subprocess.run(["brctl", "download", str(path)], timeout=timeout, capture_output=True)
    except Exception:
        pass
    try:
        with open(path, "rb") as f:
            f.read(1)
        return not is_dataless(path)
    except Exception as e:
        logger.warning("Failed to materialize dataless file %s: %s", path.name, e)
        return False


#: Marca que identifica um toco de "nao deu para extrair", para que o codigo
#: que consome o texto possa distinguir documento de aviso.
#:
#: Existe porque o toco e texto valido: ele passa por qualquer guarda de
#: `if not texto.strip()` e segue adiante como se fosse conteudo. Foi assim que
#: um deck de slides renderizados como imagem virou nota com resumo, topicos,
#: pessoas e decisoes, tudo inventado, e com `quality_score: 1.0`, porque o
#: avaliador de sintese so reprova saida com menos de 30 caracteres.
#:
#: A marca fica na PRIMEIRA linha de proposito, para que a checagem seja barata
#: e nao encontre a frase por acidente no meio de um documento que fale sobre
#: extracao de texto.
NO_TEXT_MARKER = "[no extractable text]"


#: Piso de substancia: abaixo disto nao ha documento, ha residuo.
#:
#: Nao e chute. Medido em 03/09/2026 sobre os documentos reais desta maquina,
#: 8 `.pptx`, 42 `.pdf` e 9 `.docx`:
#:
#:                            degenerados medidos      menor documento real
#:   total de caracteres      0 e 56                   4.296
#:   por slide ou pagina      0,0 e 14,0               195,1
#:
#: Os dois grupos estao separados por duas ordens de grandeza, entao qualquer
#: corte no meio do vao serve e nenhum fica ajustado a uma das bordas. 200 no
#: total e 3,5x o maior degenerado e 21x abaixo do menor documento de verdade;
#: 40 por unidade e 2,9x o maior degenerado e 4,9x abaixo do menor real.
#:
#: O piso por unidade existe para o caso que o total sozinho deixa passar: um
#: deck exportado como imagem costuma vir com zero texto, mas quando vem com
#: numero de slide ou rodape, 30 slides de 8 caracteres somam 240 e cruzam o
#: piso total sem serem documento nenhum.
#:
#: Ele so vale onde a unidade e uma pagina, isto e, slide de `.pptx` e pagina de
#: `.pdf`. Paragrafo de `.docx` nao serve de unidade: o menor arquivo real
#: medido tem 39,2 caracteres por paragrafo, abaixo do proprio piso, porque
#: paragrafo e uma unidade pequena demais para a conta significar alguma coisa.
MIN_CHARS_TOTAL = 200
MIN_CHARS_POR_PAGINA = 40


def sem_substancia(texto: str, *, paginas: int | None = None) -> bool:
    """O texto extraido e fino demais para justificar uma nota sintetizada?

    A pergunta nao e "esta vazio", e "da para escrever um resumo com topicos,
    pessoas e decisoes a partir disto". Sao coisas diferentes, e foi a distancia
    entre as duas que produziu 22 notas de invencao integral: os arquivos nao
    estavam vazios o bastante para a guarda de vazio pegar, e estavam vazios
    demais para qualquer sintese ser verdadeira.
    """
    n = len((texto or "").strip())
    if n < MIN_CHARS_TOTAL:
        return True
    if paginas and paginas > 0 and n / paginas < MIN_CHARS_POR_PAGINA:
        return True
    return False


def no_text_stub(path: Path, *, formato: str, unidade: str, total: int,
                 extraido: int = 0) -> str:
    """O registro de um arquivo do qual nao se extraiu uma linha de texto.

    Um arquivo assim nao pode virar nota sintetizada, mas tambem nao deve
    desaparecer: o PDF escaneado ja tinha esse piso e era o unico. Um deck de
    slides que sao imagens, um .docx sem paragrafo de texto e uma planilha sem
    celula preenchida caiam todos no mesmo buraco, e o que sobrava era um erro
    seco em vez de um registro de que o arquivo existe.
    """
    return (
        f"{NO_TEXT_MARKER}\n"
        f"File: {path.name}\n"
        f"Format: {formato}\n"
        f"{unidade}: {total}\n"
        f"Characters extracted: {extraido}\n"
        "There is not enough text here to summarise, so no note will be "
        "synthesised from it.\n"
        "To make it searchable, export the text yourself and drop it in as .txt."
    )


def is_no_text_stub(texto: str | None) -> bool:
    """O texto e um toco de extracao vazia, e nao um documento?"""
    return bool(texto) and texto.lstrip().startswith(NO_TEXT_MARKER)


def extract(path: Path) -> str | None:
    """
    Return extracted plain text from path.
    Raises DatalessFileError if evicted to cloud and cannot be materialized.
    Raises UnreadableFileError if parsing fails.
    Returns None if format is unsupported.
    """
    suffix = path.suffix.lower()
    _map = {
        ".md":   _text,
        ".markdown": _text,
        ".mdx":  _text,
        ".txt":  _text,
        ".text": _text,
        ".csv":  _csv_to_md,
        ".html": _html,
        ".htm":  _html,
        ".pdf":  _pdf,
        ".docx": _docx,
        ".xlsx": _xlsx,
        ".pptx": _pptx,
        ".json": _json,
    }
    fn = _map.get(suffix)
    if fn is None:
        return None

    if is_dataless(path) and not materialize_if_dataless(path):
        raise DatalessFileError(f"File {path.name} is evicted to cloud storage (dataless)")

    try:
        result = fn(path)
        return result if result is not None else ""
    except Exception as e:
        logger.warning("Extraction failed for %s: %s", path.name, e)
        raise UnreadableFileError(f"Extraction failed for {path.name}: {e}") from e



def format_label(path: Path) -> str:
    """Human-readable format name for a file path."""
    labels = {
        ".md": "Markdown", ".markdown": "Markdown",
        ".txt": "Text", ".text": "Text", ".csv": "CSV",
        ".html": "HTML", ".htm": "HTML", ".pdf": "PDF",
        ".docx": "Word", ".xlsx": "Excel", ".pptx": "PowerPoint",
        ".json": "JSON",
    }
    return labels.get(path.suffix.lower(), path.suffix.upper().lstrip("."))


# ── extractors ────────────────────────────────────────────────────────────────

def _text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _json(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        data = json.loads(raw)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        return raw


def _csv_to_md(path: Path) -> str:
    rows = []
    with open(path, encoding="utf-8", errors="replace", newline="") as fh:
        reader = _csv.reader(fh)
        for i, row in enumerate(reader):
            if not any(c.strip() for c in row):
                continue
            cells = [c.replace("|", "\\|") for c in row]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
            if i >= 500:
                rows.append("| *(table truncated at 500 rows)* |")
                break
    return "\n".join(rows)


def _html(path: Path) -> str:
    from bs4 import BeautifulSoup
    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style", "head", "nav", "footer", "aside"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _pdf(path: Path) -> str:
    import pypdf
    reader = pypdf.PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {i + 1}]\n{text.strip()}")
    texto = "\n\n".join(pages)
    if sem_substancia(texto, paginas=len(reader.pages)):
        return no_text_stub(path, formato="PDF (scanned or image-only)",
                            unidade="Pages", total=len(reader.pages),
                            extraido=len(texto.strip()))
    return texto


def _docx(path: Path) -> str:
    import docx
    doc = docx.Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    texto = "\n\n".join(parts)
    if sem_substancia(texto):
        return no_text_stub(path, formato="Word (almost no text)",
                            unidade="Paragraphs", total=len(doc.paragraphs),
                            extraido=len(texto.strip()))
    return texto


def _xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            cells = [str(c) if c is not None else "" for c in row]
            if not any(c.strip() for c in cells):
                continue
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
            if i >= 200:
                rows.append("| *(sheet truncated at 200 rows)* |")
                break
        if rows:
            parts.append(f"## {sheet_name}\n\n" + "\n".join(rows))
    n_abas = len(wb.sheetnames)
    wb.close()
    texto = "\n\n".join(parts)
    if sem_substancia(texto):
        return no_text_stub(path, formato="Excel (almost no filled cells)",
                            unidade="Sheets", total=n_abas,
                            extraido=len(texto.strip()))
    return texto


def _pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    texto = "\n\n".join(slides)
    n_slides = len(prs.slides)
    if sem_substancia(texto, paginas=n_slides):
        # Decks exportados como imagem renderizada: cada slide e um grafico. Sem
        # texto nenhum, ou com so o numero do slide, que da no mesmo.
        return no_text_stub(path, formato="PowerPoint (slides are images)",
                            unidade="Slides", total=n_slides,
                            extraido=len(texto.strip()))
    return texto


# ── Recursive split helpers (v0.3) ────────────────────────────────────────────

def extract_pages(path: Path) -> list[str]:
    """Return per-page text for PDFs. Returns [] for non-PDFs or on extraction error.

    Each element is the non-empty text of one page. Blank pages are dropped.
    Used by splitter.py to decide whether to split a PDF by page groups.
    """
    if path.suffix.lower() != ".pdf":
        return []
    try:
        import pypdf
        reader = pypdf.PdfReader(str(path))
        return [
            page.extract_text().strip()
            for page in reader.pages
            if (page.extract_text() or "").strip()
        ]
    except Exception as e:
        logger.warning("extract_pages failed for %s: %s", path.name, e)
        return []
